import os
import zipfile
from pathlib import Path

import fitz
from docx import Document
from striprtf.striprtf import rtf_to_text
from openpyxl import load_workbook
from pptx import Presentation

from langchain_core.documents import Document as LCDocument
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_chroma import Chroma


DATA_DIR = "data"
CHROMA_DIR = "chroma"

SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".rtf",
    ".txt",
    ".xlsx",
    ".pptx",
    ".zip"
}


# --------------------------------------------------
# PDF
# --------------------------------------------------

def extract_pdf(path):

    doc = fitz.open(path)

    text = []

    for page in doc:
        text.append(page.get_text())

    doc.close()

    return "\n".join(text)


# --------------------------------------------------
# DOCX
# --------------------------------------------------

def extract_docx(path):

    doc = Document(path)

    paragraphs = []

    for paragraph in doc.paragraphs:
        paragraphs.append(paragraph.text)

    return "\n".join(paragraphs)


# --------------------------------------------------
# RTF
# --------------------------------------------------

def extract_rtf(path):

    with open(
        path,
        "r",
        encoding="utf-8",
        errors="ignore"
    ) as f:

        return rtf_to_text(f.read())


# --------------------------------------------------
# TXT
# --------------------------------------------------

def extract_txt(path):

    with open(
        path,
        "r",
        encoding="utf-8",
        errors="ignore"
    ) as f:

        return f.read()


# --------------------------------------------------
# XLSX
# --------------------------------------------------

def extract_xlsx(path):

    workbook = load_workbook(
        path,
        data_only=True
    )

    text = []

    for sheet in workbook:

        text.append(f"\nSheet: {sheet.title}\n")

        for row in sheet.iter_rows(values_only=True):

            values = []

            for value in row:

                if value is None:
                    continue

                values.append(str(value))

            if values:
                text.append(" | ".join(values))

    return "\n".join(text)


# --------------------------------------------------
# PPTX
# --------------------------------------------------

def extract_pptx(path):

    prs = Presentation(path)

    text = []

    for slide in prs.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                if shape.text.strip():
                    text.append(shape.text)

    return "\n".join(text)


# --------------------------------------------------
# ZIP
# --------------------------------------------------

def extract_zip(path):

    extract_folder = os.path.splitext(path)[0]

    os.makedirs(extract_folder, exist_ok=True)

    with zipfile.ZipFile(path, "r") as zip_ref:
        zip_ref.extractall(extract_folder)

    print(f"📦 Extracted: {path}")

    return extract_folder


# --------------------------------------------------
# Dispatcher
# --------------------------------------------------

def extract_text(path):

    suffix = Path(path).suffix.lower()

    if suffix == ".pdf":
        return extract_pdf(path)

    if suffix == ".docx":
        return extract_docx(path)

    if suffix == ".rtf":
        return extract_rtf(path)

    if suffix == ".txt":
        return extract_txt(path)

    if suffix == ".xlsx":
        return extract_xlsx(path)

    if suffix == ".pptx":
        return extract_pptx(path)

    if suffix == ".zip":

        extract_zip(path)

        return None

    return None


def process_directory(folder, splitter, seen_chunks):

    documents = []

    for root, _, files in os.walk(folder):

        for file in files:

            path = os.path.join(root, file)

            suffix = Path(path).suffix.lower()

            if suffix not in SUPPORTED_EXTENSIONS:
                continue

            print(f"Processing: {path}")

            try:

                # ZIP
                if suffix == ".zip":

                    extracted_folder = extract_zip(path)

                    documents.extend(
                        process_directory(
                            extracted_folder,
                            splitter,
                            seen_chunks
                        )
                    )

                    continue

                text = extract_text(path)

                if text is None:
                    continue

                if not text.strip():
                    continue

                chunks = splitter.split_text(text)

                for chunk in chunks:

                    key = chunk.strip()

                    if len(key) < 30:
                        continue

                    if key in seen_chunks:
                        continue

                    seen_chunks.add(key)

                    company = (
                        os.path.basename(root)
                        if os.path.abspath(root) != os.path.abspath(DATA_DIR)
                        else os.path.splitext(file)[0]
                    )

                    documents.append(
                        LCDocument(
                            page_content=chunk,
                            metadata={
                                "source": file,
                                "filename": file,
                                "path": os.path.relpath(path, DATA_DIR),
                                "extension": suffix,
                                "company": company,
                            },
                        )
                    )

            except Exception as e:

                print(f"❌ {path}")

                print(e)

    return documents


# --------------------------------------------------
# Walk all files recursively
# --------------------------------------------------

def collect_documents():

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200
    )

    seen_chunks = set()

    return process_directory(
        DATA_DIR,
        splitter,
        seen_chunks
    )


# --------------------------------------------------
# Build Chroma
# --------------------------------------------------

def build_chroma(documents):

    print("\nLoading embedding model...")

    embedding_model = HuggingFaceEmbeddings(

        model_name="sentence-transformers/all-MiniLM-L6-v2"

    )

    if os.path.exists(CHROMA_DIR):

        import shutil

        shutil.rmtree(CHROMA_DIR)

    print("\nCreating ChromaDB...\n")

    batch_size = 100

    db = Chroma(

        persist_directory=CHROMA_DIR,

        embedding_function=embedding_model

    )

    for i in range(0, len(documents), batch_size):

        batch = documents[i:i + batch_size]

        db.add_documents(batch)

        print(

            f"Indexed "

            f"{min(i + batch_size, len(documents))}"

            f"/"

            f"{len(documents)}"

        )

    print("\n✅ Chroma Database Created")


# --------------------------------------------------
# Main
# --------------------------------------------------

def main():

    print("\nScanning files...\n")

    documents = collect_documents()

    print(f"\nDocuments: {len(documents)} chunks")

    build_chroma(documents)

    print("\n🎉 Ingestion Complete")


if __name__ == "__main__":

    main()